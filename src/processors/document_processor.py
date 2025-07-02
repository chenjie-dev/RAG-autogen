import os
from typing import List, Dict, Any
from pathlib import Path
import docx
import markdown
from bs4 import BeautifulSoup
from pptx import Presentation
import json

# DocLing imports
from docling.backend.docling_parse_v2_backend import DoclingParseV2DocumentBackend
from docling.document_converter import DocumentConverter, FormatOption
from docling.datamodel.pipeline_options import PdfPipelineOptions, TableFormerMode, EasyOcrOptions
from docling.datamodel.base_models import InputFormat
from docling.pipeline.standard_pdf_pipeline import StandardPdfPipeline
from docling.datamodel.base_models import ConversionStatus
from docling.datamodel.document import ConversionResult, InputDocument

# 导入日志模块
from utils.logger import logger

class DocumentProcessor:
    """文档处理器，支持多种文件格式，使用DocLing处理PDF"""
    
    @staticmethod
    def _create_document_converter():
        """创建文档转换器，使用DocLingParseV2DocumentBackend"""
        pipeline_options = PdfPipelineOptions()
        pipeline_options.do_ocr = True
        ocr_options = EasyOcrOptions(lang=['en'], force_full_page_ocr=False)
        pipeline_options.ocr_options = ocr_options
        pipeline_options.do_table_structure = True
        pipeline_options.table_structure_options.do_cell_matching = True
        pipeline_options.table_structure_options.mode = TableFormerMode.ACCURATE
        
        format_options = {
            InputFormat.PDF: FormatOption(
                pipeline_cls=StandardPdfPipeline,
                pipeline_options=pipeline_options,
                backend=DoclingParseV2DocumentBackend
            )
        }
        
        return DocumentConverter(format_options=format_options)

    @staticmethod
    def process_pdf_with_pages(file_path: str) -> Dict[str, Any]:
        """处理PDF文件并保持页面信息，返回包含页面和chunks的结构化数据"""
        try:
            logger.info(f"开始处理PDF文件并保持页面信息: {os.path.basename(file_path)}")
            
            # 创建文档转换器
            doc_converter = DocumentProcessor._create_document_converter()
            
            # 转换文档
            input_path = Path(file_path)
            conv_results = doc_converter.convert_all(source=[input_path])
            
            for conv_result in conv_results:
                if conv_result.status == ConversionStatus.SUCCESS:
                    # 获取文档数据
                    data = conv_result.document.export_to_dict()
                    
                    # 处理文档内容，保持页面信息
                    processed_data = DocumentProcessor._process_document_with_pages(data)
                    
                    logger.info(f"PDF文件处理完成，共 {len(processed_data.get('pages', []))} 页，{len(processed_data.get('chunks', []))} 个chunks")
                    return processed_data
                    
                else:
                    logger.error(f"PDF文档转换失败: {conv_result.status}")
                    return {"pages": [], "chunks": []}
                    
        except Exception as e:
            logger.error(f"处理PDF文件时出错: {str(e)}")
            import traceback
            logger.error(f"错误详情: {traceback.format_exc()}")
            return {"pages": [], "chunks": []}

    @staticmethod
    def _process_document_with_pages(data: Dict[str, Any]) -> Dict[str, Any]:
        """处理文档数据，提取页面和chunks信息"""
        result = {
            "pages": [],
            "chunks": []
        }
        
        # 添加调试日志，查看数据结构
        logger.info(f"DocLing返回的数据结构: {list(data.keys())}")
        
        # 处理页面信息 - 尝试多种可能的数据结构
        page_texts = {}
        
        # 方法1: 直接从texts字段提取，按页面分组
        if 'texts' in data:
            logger.info(f"找到texts字段，包含 {len(data['texts'])} 个文本项")
            for text_item in data['texts']:
                if 'text' in text_item and text_item['text'].strip():
                    # 获取页面信息
                    page_num = 0
                    if 'prov' in text_item and text_item['prov']:
                        prov = text_item['prov'][0]
                        page_num = prov.get('page_no', 0)
                    
                    if page_num not in page_texts:
                        page_texts[page_num] = ""
                    page_texts[page_num] += text_item['text'].strip() + "\n"
        
        # 方法2: 如果content字段存在，按页面处理
        elif 'content' in data:
            logger.info(f"找到content字段，包含 {len(data['content'])} 个内容项")
            for page_data in data['content']:
                page_num = page_data.get('page', 0)
                page_text = ""
                
                if 'texts' in page_data:
                    for text_item in page_data['texts']:
                        if 'text' in text_item and text_item['text'].strip():
                            page_text += text_item['text'].strip() + "\n"
                
                if page_text.strip():
                    page_texts[page_num] = page_text.strip()
        
        # 方法3: 如果pages字段存在
        elif 'pages' in data:
            logger.info(f"找到pages字段，包含 {len(data['pages'])} 页")
            for page_data in data['pages']:
                page_num = page_data.get('page_no', 0)
                page_text = ""
                
                if 'texts' in page_data:
                    for text_item in page_data['texts']:
                        if 'text' in text_item and text_item['text'].strip():
                            page_text += text_item['text'].strip() + "\n"
                
                if page_text.strip():
                    page_texts[page_num] = page_text.strip()
        
        # 如果没有找到任何文本，尝试直接处理整个文档
        if not page_texts:
            logger.warning("未找到页面结构，尝试直接处理整个文档")
            # 将所有文本合并为一个页面
            all_text = ""
            if 'texts' in data:
                for text_item in data['texts']:
                    if 'text' in text_item and text_item['text'].strip():
                        all_text += text_item['text'].strip() + "\n"
            
            if all_text.strip():
                page_texts[0] = all_text.strip()
        
        # 构建结果
        for page_num, page_text in sorted(page_texts.items()):
            if page_text.strip():
                result["pages"].append({
                    "page": page_num,
                    "text": page_text.strip()
                })
                
                # 分割页面文本为chunks
                chunks = DocumentProcessor._split_page_text(page_text.strip(), page_num)
                result["chunks"].extend(chunks)
        
        logger.info(f"处理完成，共提取 {len(result['pages'])} 页，{len(result['chunks'])} 个chunks")
        return result

    @staticmethod
    def _split_page_text(text: str, page_num: int, chunk_size: int = 500, chunk_overlap: int = 50) -> List[Dict[str, Any]]:
        """将页面文本分割为chunks，保持页面信息"""
        chunks = []
        
        # 简单的文本分割（可以后续优化为更智能的分割）
        words = text.split()
        current_chunk = []
        current_length = 0
        
        for word in words:
            current_chunk.append(word)
            current_length += len(word) + 1  # +1 for space
            
            if current_length >= chunk_size:
                chunk_text = " ".join(current_chunk)
                chunks.append({
                    "page": page_num,
                    "text": chunk_text,
                    "length": len(chunk_text)
                })
                
                # 保留重叠部分
                overlap_words = current_chunk[-chunk_overlap//10:] if chunk_overlap > 0 else []
                current_chunk = overlap_words
                current_length = sum(len(word) + 1 for word in overlap_words)
        
        # 添加最后一个chunk
        if current_chunk:
            chunk_text = " ".join(current_chunk)
            chunks.append({
                "page": page_num,
                "text": chunk_text,
                "length": len(chunk_text)
            })
        
        return chunks

    @staticmethod
    def extract_text_from_pdf(file_path: str) -> List[str]:
        """从PDF文件中提取文本，使用DocLingParseV2DocumentBackend"""
        texts = []
        try:
            logger.info(f"开始处理PDF文件: {os.path.basename(file_path)}")
            
            # 创建文档转换器
            doc_converter = DocumentProcessor._create_document_converter()
            
            # 转换文档
            input_path = Path(file_path)
            conv_results = doc_converter.convert_all(source=[input_path])
            
            for conv_result in conv_results:
                if conv_result.status == ConversionStatus.SUCCESS:
                    # 获取文档数据
                    data = conv_result.document.export_to_dict()
                    
                    # 处理文本内容
                    if 'texts' in data:
                        for text_item in data['texts']:
                            if 'text' in text_item and text_item['text'].strip():
                                texts.append(text_item['text'].strip())
                    
                    # 处理表格内容
                    if 'tables' in data:
                        for table_item in data['tables']:
                            if 'data' in table_item and 'grid' in table_item['data']:
                                for row in table_item['data']['grid']:
                                    row_texts = []
                                    for cell in row:
                                        if 'text' in cell and cell['text'].strip():
                                            row_texts.append(cell['text'].strip())
                                    if row_texts:
                                        texts.append(" | ".join(row_texts))
                    
                    logger.info(f"使用DocLing处理PDF文件，提取了 {len(texts)} 个文本块")
                    
                    # 显示提取的文本示例
                    if texts:
                        logger.info("提取的文本示例:")
                        for i, text in enumerate(texts[:3]):
                            logger.info(f"  {i+1}: {text[:100]}...")
                    
                else:
                    logger.error(f"PDF文档转换失败: {conv_result.status}")
                    
        except Exception as e:
            logger.error(f"使用DocLing处理PDF文件时出错: {str(e)}")
            import traceback
            logger.error(f"错误详情: {traceback.format_exc()}")
            
        return texts

    @staticmethod
    def extract_text_from_pdf_with_layout(file_path: str) -> List[dict]:
        """从PDF文件中提取文本和布局信息，使用DocLingParseV2DocumentBackend"""
        results = []
        try:
            logger.info(f"开始处理PDF文件布局: {os.path.basename(file_path)}")
            
            # 创建文档转换器
            doc_converter = DocumentProcessor._create_document_converter()
            
            # 转换文档
            input_path = Path(file_path)
            conv_results = doc_converter.convert_all(source=[input_path])
            
            for conv_result in conv_results:
                if conv_result.status == ConversionStatus.SUCCESS:
                    # 获取文档数据
                    data = conv_result.document.export_to_dict()
                    
                    # 处理文本内容
                    if 'texts' in data:
                        for text_item in data['texts']:
                            if 'text' in text_item and text_item['text'].strip():
                                item_data = {
                                    'text': text_item['text'].strip(),
                                    'type': 'text'
                                }
                                
                                # 添加位置信息
                                if 'prov' in text_item and text_item['prov']:
                                    prov = text_item['prov'][0]
                                    item_data['page'] = prov.get('page_no')
                                    if 'bbox' in prov:
                                        item_data['bbox'] = {
                                            'left': prov['bbox'].get('l'),
                                            'top': prov['bbox'].get('t'),
                                            'right': prov['bbox'].get('r'),
                                            'bottom': prov['bbox'].get('b')
                                        }
                                
                                # 添加字体信息
                                if 'font_size' in text_item:
                                    item_data['font_size'] = text_item['font_size']
                                if 'font_name' in text_item:
                                    item_data['font_name'] = text_item['font_name']
                                
                                results.append(item_data)
                    
                    logger.info(f"使用DocLing处理PDF文件布局，提取了 {len(results)} 个文本块")
                    
                else:
                    logger.error(f"PDF文档转换失败: {conv_result.status}")
                    
        except Exception as e:
            logger.error(f"使用DocLing处理PDF文件布局时出错: {str(e)}")
            import traceback
            logger.error(f"错误详情: {traceback.format_exc()}")
            
        return results

    @staticmethod
    def extract_text_from_pdf_in_region(file_path: str, bbox: dict, page_index: int = 0) -> str:
        """从PDF文件的指定区域提取文本"""
        try:
            logger.info(f"从PDF区域提取文本: {os.path.basename(file_path)}")
            
            # 创建文档转换器
            doc_converter = DocumentProcessor._create_document_converter()
            
            # 转换文档
            input_path = Path(file_path)
            conv_results = doc_converter.convert_all(source=[input_path])
            
            for conv_result in conv_results:
                if conv_result.status == ConversionStatus.SUCCESS:
                    # 获取文档数据
                    data = conv_result.document.export_to_dict()
                    
                    # 查找指定页面和区域的文本
                    if 'texts' in data:
                        for text_item in data['texts']:
                            if 'prov' in text_item and text_item['prov']:
                                prov = text_item['prov'][0]
                                if prov.get('page_no') == page_index:
                                    # 检查文本是否在指定区域内
                                    if 'bbox' in prov:
                                        text_bbox = prov['bbox']
                                        # 简单的区域重叠检查
                                        if (text_bbox['l'] >= bbox['left'] and 
                                            text_bbox['t'] >= bbox['top'] and
                                            text_bbox['r'] <= bbox['right'] and
                                            text_bbox['b'] <= bbox['bottom']):
                                            if 'text' in text_item and text_item['text'].strip():
                                                return text_item['text'].strip()
                    
                else:
                    logger.error(f"PDF文档转换失败: {conv_result.status}")
                    
        except Exception as e:
            logger.error(f"从PDF区域提取文本时出错: {str(e)}")
            
        return ""

    @staticmethod
    def extract_text_from_docx(file_path: str) -> List[str]:
        """从Word文档中提取文本，使用DocLing"""
        texts = []
        try:
            logger.info(f"开始处理Word文档: {os.path.basename(file_path)}")
            
            # 使用DocLing的MsWordDocumentBackend
            from docling.backend.msword_backend import MsWordDocumentBackend
            from docling.datamodel.document import DoclingDocument, TextItem, SectionHeaderItem
            
            # 创建输入文档
            in_doc = InputDocument(
                path_or_stream=Path(file_path),
                format=InputFormat.DOCX,
                backend=MsWordDocumentBackend,
            )
            
            # 创建后端并转换文档
            backend = MsWordDocumentBackend(
                in_doc=in_doc,
                path_or_stream=Path(file_path),
            )
            doc: DoclingDocument = backend.convert()
            
            logger.info(f"DocLing文档转换成功，开始提取文本...")
            
            # 遍历文档中的所有项目
            item_count = 0
            for item, _ in doc.iterate_items():
                item_count += 1
                
                if isinstance(item, TextItem):
                    # 处理普通文本
                    if item.text.strip():
                        texts.append(item.text.strip())
                elif isinstance(item, SectionHeaderItem):
                    # 处理标题，添加标题级别信息
                    header_text = f"[标题{item.level}] {item.text.strip()}"
                    if header_text:
                        texts.append(header_text)
                else:
                    # 处理其他类型的项目
                    if hasattr(item, 'text') and item.text.strip():
                        texts.append(item.text.strip())
            
            logger.info(f"使用DocLing处理Word文档，遍历了 {item_count} 个项目，提取了 {len(texts)} 个文本块")
            
            # 验证提取结果
            if not texts:
                logger.warning("DocLing没有提取到任何文本")
                raise Exception("DocLing提取结果为空")
            
            # 显示提取的文本示例
            if texts:
                logger.info("提取的文本示例:")
                for i, text in enumerate(texts[:3]):
                    logger.info(f"  {i+1}: {text[:100]}...")
            
        except Exception as e:
            logger.error(f"使用DocLing处理Word文档时出错: {str(e)}")
            import traceback
            logger.error(f"错误详情: {traceback.format_exc()}")
            texts = []
        
        return texts

    @staticmethod
    def extract_text_from_docx_with_structure(file_path: str) -> List[dict]:
        """从Word文档中提取文本和结构信息，使用DocLing"""
        results = []
        try:
            # 使用DocLing的MsWordDocumentBackend
            from docling.backend.msword_backend import MsWordDocumentBackend
            from docling.datamodel.document import DoclingDocument, TextItem, SectionHeaderItem
            
            # 创建输入文档
            in_doc = InputDocument(
                path_or_stream=Path(file_path),
                format=InputFormat.DOCX,
                backend=MsWordDocumentBackend,
            )
            
            # 创建后端并转换文档
            backend = MsWordDocumentBackend(
                in_doc=in_doc,
                path_or_stream=Path(file_path),
            )
            doc: DoclingDocument = backend.convert()
            
            # 遍历文档中的所有项目
            for item, _ in doc.iterate_items():
                item_data = {
                    'text': '',
                    'type': 'unknown',
                    'level': None,
                    'metadata': {}
                }
                
                if isinstance(item, TextItem):
                    item_data['text'] = item.text.strip()
                    item_data['type'] = 'text'
                    if hasattr(item, 'font_size'):
                        item_data['metadata']['font_size'] = getattr(item, 'font_size', None)
                    if hasattr(item, 'font_name'):
                        item_data['metadata']['font_name'] = getattr(item, 'font_name', None)
                        
                elif isinstance(item, SectionHeaderItem):
                    item_data['text'] = item.text.strip()
                    item_data['type'] = 'header'
                    item_data['level'] = item.level
                    
                else:
                    # 处理其他类型的项目
                    if hasattr(item, 'text'):
                        item_data['text'] = item.text.strip()
                    item_data['type'] = type(item).__name__
                
                # 只添加有文本内容的项目
                if item_data['text']:
                    results.append(item_data)
            
            logger.info(f"使用DocLing处理Word文档，提取了 {len(results)} 个结构化文本块")
            
        except Exception as e:
            logger.error(f"使用DocLing处理Word文档结构时出错: {str(e)}")
            import traceback
            logger.error(f"错误详情: {traceback.format_exc()}")
        return results

    @staticmethod
    def extract_text_from_markdown(file_path: str) -> List[str]:
        """从Markdown文件中提取文本"""
        texts = []
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                md_content = file.read()
                # 转换为HTML
                html = markdown.markdown(md_content)
                # 使用BeautifulSoup提取文本
                soup = BeautifulSoup(html, 'html.parser')
                # 获取所有段落
                paragraphs = soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6'])
                for p in paragraphs:
                    text = p.get_text().strip()
                    if text:
                        texts.append(text)
        except Exception as e:
            logger.error(f"处理Markdown文件时出错: {str(e)}")
        return texts

    @staticmethod
    def extract_text_from_pptx(file_path: str) -> List[str]:
        """从PowerPoint文件中提取文本"""
        texts = []
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_texts.append(text)
                if slide_texts:
                    texts.append(" | ".join(slide_texts))  # 将同一页的文本合并
        except Exception as e:
            logger.error(f"处理PowerPoint文件时出错: {str(e)}")
        return texts

    @staticmethod
    def extract_text_from_txt(file_path: str) -> List[str]:
        """从文本文件中提取文本"""
        texts = []
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
                # 按段落分割
                paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
                texts.extend(paragraphs)
        except Exception as e:
            logger.error(f"处理文本文件时出错: {str(e)}")
        return texts

    @staticmethod
    def process_file(file_path: str) -> List[str]:
        """处理各种格式的文件并提取文本"""
        file_ext = os.path.splitext(file_path)[1].lower()
        
        processors = {
            '.pdf': DocumentProcessor.extract_text_from_pdf,
            '.docx': DocumentProcessor.extract_text_from_docx,
            '.md': DocumentProcessor.extract_text_from_markdown,
            '.pptx': DocumentProcessor.extract_text_from_pptx,
            '.txt': DocumentProcessor.extract_text_from_txt
        }
        
        if file_ext not in processors:
            raise ValueError(f"不支持的文件格式: {file_ext}")
            
        return processors[file_ext](file_path) 